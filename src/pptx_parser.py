import logging
import os
import re
from datetime import datetime

from pptx import Presentation

from utilities import (format_date, get_dir_for_file_path, get_reference,
                       sanitize_filename)

date_pattern = re.compile(r'\d{4}_\d{2}_\d{2}')
reference_pattern = re.compile(r'^[a-zA-Z0-9]+(\.[a-zA-Z0-9]+)+$')


def escape_markdown_characters(text):
    markdown_chars = ['*', '_', '`', '[', ']',
                      '(', ')', '#', '!']
    for char in markdown_chars:
        text = text.replace(char, f"\\{char}")
    return text


def extract_rich_text_from_paragraph(paragraph):
    rich_text = ""
    for run in paragraph.runs:
        text = escape_markdown_characters(run.text)

        if run.font.bold:
            text = f"**{text}**"
        if run.font.italic:
            text = f"_{text}_"
        if run.font.underline:
            text = f"<u>{text}</u>"
        try:
            if run.hyperlink.address != None:
                text = f"[{text}]({run.hyperlink.address})"
        except:
            pass
        rich_text += text
    return rich_text


def parse_pptx(file_path, base_dir="../output_markdown_files"):
    logging.debug(f'Processing file: {file_path}')

    pres = Presentation(file_path)
    slides_data = []

    output_dir = os.path.join(base_dir, get_dir_for_file_path(file_path))
    os.makedirs(output_dir, exist_ok=True)

    # insure output directory exists

    for index, slide in enumerate(pres.slides):
        logging.debug(f'    > Slide {index + 1}')
        slide_data = {}

        # Extract title, content, and dates
        title = slide.shapes.title.text.strip(
        ) if slide.shapes.title else f"Slide {index + 1}"
        content = ""
        date_pattern = re.compile(r'\d{4}_\d{2}_\d{2}')
        creation_date = ""
        update_dates = []
        images = []

        img_idx = 1
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture shape type
                image = shape.image
                if image.filename == None:
                    image_name = f"{title}-{img_idx}.png"
                else:
                    image_name = f"{title}-{img_idx}-{image.filename}"

                image_name = sanitize_filename(image_name)
                img_path = os.path.join(
                    output_dir, "img", image_name)
                img_idx += 1

                os.makedirs(os.path.dirname(img_path), exist_ok=True)

                with open(img_path, "wb") as img_file:
                    img_file.write(image.blob)

                content += f"\n![{title} Image {image_name}](img/{image_name})\n"

            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()

                if reference_pattern.match(text):
                    slide_data["reference"] = get_reference(text)
                if date_pattern.match(text):
                    date = format_date(text)
                    if date != None:
                        if not creation_date:
                            creation_date = format_date(text)
                        else:
                            update_dates.append(format_date(text))
                else:
                    content += extract_rich_text_from_paragraph(
                        paragraph) + "\n"

        slide_data["title"] = title
        slide_data["content"] = content
        slide_data["creation_date"] = creation_date
        slide_data["update_dates"] = update_dates

        slides_data.append(slide_data)

    core_properties = [x for x in dir(
        pres.core_properties) if not x.startswith("_")]
    meta = {}
    for key in core_properties:
        if getattr(pres.core_properties, key):
            if type(getattr(pres.core_properties, key)) == str:
                meta[key] = getattr(pres.core_properties, key)
            elif type(getattr(pres.core_properties, key)) == datetime:
                meta[key] = getattr(pres.core_properties, key).strftime(
                    "%d/%m/%Y %H:%M:%S")

    return [slides_data, meta, output_dir]
